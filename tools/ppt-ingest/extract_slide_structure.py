#!/usr/bin/env python3
"""
extract_slide_structure.py
python-pptx를 사용하여 PPTX 내 각 슬라이드의 도형 구조, 텍스트, 좌표 및 HTML을 추출합니다.
(Deterministic Extraction)
"""

import argparse
import html
import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


WORKBENCH_ROOT = Path(__file__).resolve().parents[2]
ASSET_CURATOR_DIR = WORKBENCH_ROOT / "tools" / "asset-curator"
if str(ASSET_CURATOR_DIR) not in sys.path:
    sys.path.insert(0, str(ASSET_CURATOR_DIR))
from asset_curator import describe_source, extract_index_slides


def emu_to_px(v: int) -> int:
    """EMU 단위를 픽셀(96 DPI 기준 약 9525 EMU/px)로 변환"""
    return int(v / 9525)


def iter_shapes(shapes):
    """그룹 도형을 재귀적으로 순회하여 내부 도형들을 yield"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def clean_text(text: str) -> str:
    if not text:
        return ""
    return text.strip().replace("\x0b", "\n").replace("\r", "\n")


def extract_year_from_name(filename: str) -> int:
    """파일명에서 202x 형태의 연도 추출 (기본값: 2026)"""
    match = re.search(r"(202\d)", filename)
    if match:
        return int(match.group(1))
    return 2026


def build_slide_html(slide_w_px: int, slide_h_px: int, elements: list[str]) -> str:
    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<style>
body {{
    margin: 0;
    background: #f0f2f5;
}}
.slide {{
    position: relative;
    width: {slide_w_px}px;
    height: {slide_h_px}px;
    background: white;
    overflow: hidden;
    font-family: Arial, "Malgun Gothic", sans-serif;
    box-shadow: 0 4px 12px rgba(0,0,0,0.08);
}}
.shape {{
    position: absolute;
    box-sizing: border-box;
    font-size: 14px;
    line-height: 1.35;
    white-space: normal;
    overflow: hidden;
    padding: 3px;
    word-break: keep-all;
}}
</style>
</head>
<body>
<section class="slide">
{''.join(elements)}
</section>
</body>
</html>"""


def _part_id(part) -> str:
    partname = getattr(part, "partname", "")
    return str(partname).rsplit("/", 1)[-1] if partname else ""


def _extract_potx_slides(potx_file: Path, html_out_dir: str | None = None) -> list[dict]:
    package = describe_source(str(potx_file))
    indexed_slides = extract_index_slides(str(potx_file))
    slide_w_px = int(package["slide_size"]["px"]["width"])
    slide_h_px = int(package["slide_size"]["px"]["height"])
    year = extract_year_from_name(potx_file.name)
    out_html_path = Path(html_out_dir).resolve() if html_out_dir else None
    if out_html_path:
        out_html_path.mkdir(parents=True, exist_ok=True)
    result = []
    for indexed in indexed_slides:
        elements = []
        texts = []
        for box in indexed.get("text_boxes", []):
            text = clean_text(box.get("text", ""))
            if not text:
                continue
            texts.append(text)
            bounds = box.get("bounds", {})
            left = int(float(bounds.get("x", 0)) * slide_w_px)
            top = int(float(bounds.get("y", 0)) * slide_h_px)
            width = int(float(bounds.get("w", 0)) * slide_w_px)
            height = int(float(bounds.get("h", 0)) * slide_h_px)
            elements.append(f"""
        <div class="shape" style="left:{left}px; top:{top}px; width:{width}px; height:{height}px;">
            {html.escape(text).replace(chr(10), "<br>")}
        </div>""")
        slide_no = indexed["slide_no"]
        html_file_name = f"slide_{slide_no:02d}.html"
        html_content = build_slide_html(slide_w_px, slide_h_px, elements)
        if out_html_path:
            (out_html_path / html_file_name).write_text(html_content, encoding="utf-8")
        result.append({
            "slide_no": slide_no,
            "source_pptx": potx_file.name,
            "source_type": "potx",
            "year": year,
            "title": indexed.get("title", "").strip(),
            "raw_text": "\n".join(texts),
            "text_count": len(texts),
            "html_file_name": html_file_name,
            "html_content": html_content,
            "layout_id": indexed.get("layout_id", ""),
            "master_id": indexed.get("master_id", ""),
        })
    return result


def extract_slides(pptx_path: str, html_out_dir: str | None = None) -> list[dict]:
    pptx_file = Path(pptx_path).resolve()
    if not pptx_file.exists():
        raise FileNotFoundError(f"PowerPoint source not found: {pptx_file}")
    if pptx_file.suffix.lower() not in {".pptx", ".potx"}:
        raise ValueError("Only .pptx and .potx sources are supported.")

    if pptx_file.suffix.lower() == ".potx":
        return _extract_potx_slides(pptx_file, html_out_dir)

    prs = Presentation(str(pptx_file))
    slide_w_px = emu_to_px(prs.slide_width)
    slide_h_px = emu_to_px(prs.slide_height)
    year = extract_year_from_name(pptx_file.name)

    if html_out_dir:
        out_html_path = Path(html_out_dir).resolve()
        out_html_path.mkdir(parents=True, exist_ok=True)

    extracted_slides = []

    for i, slide in enumerate(prs.slides, start=1):
        elements = []
        texts = []
        candidate_titles = []

        for shape in iter_shapes(slide.shapes):
            left = emu_to_px(shape.left)
            top = emu_to_px(shape.top)
            width = emu_to_px(shape.width)
            height = emu_to_px(shape.height)

            if left + width < 0 or top + height < 0:
                continue
            if left > slide_w_px or top > slide_h_px:
                continue
            if not hasattr(shape, "text"):
                continue

            text = clean_text(shape.text)
            if not text:
                continue

            texts.append(text)

            # 상단(top <= 높이의 25%)에 위치한 텍스트를 제목 후보로 수집
            if top <= slide_h_px * 0.25 and len(text) < 100:
                candidate_titles.append((top, text))

            text_html = html.escape(text).replace("\n", "<br>")
            elements.append(f"""
        <div class="shape" style="left:{left}px; top:{top}px; width:{width}px; height:{height}px;">
            {text_html}
        </div>""")

        # 타이틀 결정: shape.title 이 있으면 우선, 없으면 상단 가장 위쪽 텍스트
        slide_title = ""
        if slide.shapes.title and hasattr(slide.shapes.title, "text") and clean_text(slide.shapes.title.text):
            slide_title = clean_text(slide.shapes.title.text).split("\n")[0]
        elif candidate_titles:
            candidate_titles.sort(key=lambda x: x[0])
            slide_title = candidate_titles[0][1].split("\n")[0]
        else:
            slide_title = texts[0].split("\n")[0] if texts else f"슬라이드 {i}"

        slide_html_content = build_slide_html(slide_w_px, slide_h_px, elements)
        html_file_name = f"slide_{i:02d}.html"

        if html_out_dir:
            (out_html_path / html_file_name).write_text(slide_html_content, encoding="utf-8")

        extracted_slides.append({
            "slide_no": i,
            "source_pptx": pptx_file.name,
            "source_type": "pptx",
            "year": year,
            "title": slide_title.strip(),
            "raw_text": "\n".join(texts),
            "text_count": len(texts),
            "html_file_name": html_file_name,
            "html_content": slide_html_content,
            "layout_id": _part_id(getattr(slide, "slide_layout", None).part) if getattr(slide, "slide_layout", None) else "",
            "master_id": _part_id(getattr(getattr(slide, "slide_layout", None), "slide_master", None).part) if getattr(getattr(slide, "slide_layout", None), "slide_master", None) else "",
        })

    return extracted_slides


def main():
    parser = argparse.ArgumentParser(description="Extract slide texts, structure and HTML from PPTX.")
    parser.add_argument("--pptx", required=True, help="Path to .pptx file")
    parser.add_argument("--output-html-dir", help="Directory to save generated slide HTMLs")
    parser.add_argument("--output-json", help="Path to save structured JSON manifest")

    args = parser.parse_args()
    slides = extract_slides(args.pptx, args.output_html_dir)

    print(f"[Structure Extractor] Extracted {len(slides)} slides from {args.pptx}")
    for s in slides[:3]:
        print(f" - Slide {s['slide_no']}: {s['title']} ({s['text_count']} text blocks)")

    if args.output_json:
        # html_content는 용량이 클 수 있으므로 manifest에서는 제외하거나 간략화
        manifest = [{k: v for k, v in s.items() if k != "html_content"} for s in slides]
        Path(args.output_json).write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[Structure Extractor] Saved manifest to {args.output_json}")


if __name__ == "__main__":
    main()
