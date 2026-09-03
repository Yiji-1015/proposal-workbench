"""PPTX 구조·미디어 자산 감사 CLI. 인제스트 대상 원본을 도형/이미지 단위로 점검한다."""

from __future__ import annotations

import argparse
import hashlib
import json
from collections import Counter, defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def shape_kind(shape) -> str:
    tag = shape.element.tag.rsplit("}", 1)[-1]
    if tag == "grpSp":
        return "group"
    if tag == "cxnSp":
        return "connector"
    if tag == "pic":
        return "picture"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    return getattr(shape.shape_type, "name", str(shape.shape_type)).lower()


def child_count(shape) -> int:
    if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
        return 0
    return sum(1 + child_count(child) for child in shape.shapes)


def text_snippet(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join(shape.text.split())[:120]


def analyze_shapes(shapes, slide_width: int, slide_height: int, media_by_hash, stats, depth=0):
    records = []
    for shape in shapes:
        kind = shape_kind(shape)
        stats["all_shapes"] += 1
        stats["kinds"][kind] += 1
        stats["max_depth"] = max(stats["max_depth"], depth)
        if kind == "group":
            stats["groups"] += 1
        elif kind == "picture":
            stats["pictures"] += 1
        elif kind == "connector":
            stats["connectors"] += 1
        elif kind == "table":
            stats["tables"] += 1
        elif kind == "chart":
            stats["charts"] += 1
        else:
            stats["native_primitives"] += 1

        text = text_snippet(shape)
        if text:
            stats["text_shapes"] += 1
            stats["text_chars"] += len(text)

        width = max(0, int(getattr(shape, "width", 0)))
        height = max(0, int(getattr(shape, "height", 0)))
        area_ratio = (width * height) / (slide_width * slide_height)
        if depth == 0 and kind == "picture":
            stats["top_picture_area_ratio"] += area_ratio

        record = {
            "name": shape.name,
            "kind": kind,
            "depth": depth,
            "bounds": {
                "x": round(int(shape.left) / slide_width, 4),
                "y": round(int(shape.top) / slide_height, 4),
                "w": round(width / slide_width, 4),
                "h": round(height / slide_height, 4),
            },
            "text": text,
            "descendants": child_count(shape),
        }
        assert kind == "group" or record["descendants"] == 0

        if kind == "picture":
            blob = shape.image.blob
            digest = sha256(blob)
            record["media_sha256"] = digest
            record["media_file"] = media_by_hash.get(digest, {}).get("name", shape.image.filename)
            stats["media_hashes"].append(digest)

        if kind == "group":
            record["children"] = analyze_shapes(
                shape.shapes, slide_width, slide_height, media_by_hash, stats, depth + 1
            )
        records.append(record)
    return records


def summarize_groups(records):
    summaries = []

    def flatten(children):
        nodes = []
        for child in children:
            nodes.append(child)
            nodes.extend(flatten(child.get("children", [])))
        return nodes

    def visit(nodes):
        for node in nodes:
            if node["kind"] == "group":
                descendants = flatten(node.get("children", []))
                area = node["bounds"]["w"] * node["bounds"]["h"]
                pictures = sum(child["kind"] == "picture" for child in descendants)
                text_chars = sum(len(child.get("text", "")) for child in descendants)
                if 1 <= len(descendants) <= 40 and 0.001 <= area <= 0.2 and pictures == 0 and text_chars <= 24:
                    candidate_kind = "icon"
                elif 3 <= len(descendants) <= 200 and 0.05 <= area <= 0.9:
                    candidate_kind = "diagram"
                else:
                    candidate_kind = "other"
                summaries.append(
                    {
                        "name": node["name"],
                        "depth": node["depth"],
                        "descendants": len(descendants),
                        "pictures": pictures,
                        "text_chars": text_chars,
                        "area_ratio": round(area, 4),
                        "candidate_kind": candidate_kind,
                        "bounds": node["bounds"],
                    }
                )
            visit(node.get("children", []))

    visit(records)
    return summaries


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("media_dir", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    media = {}
    for path in sorted(args.media_dir.iterdir()):
        if not path.is_file():
            continue
        data = path.read_bytes()
        info = {
            "name": path.name,
            "bytes": len(data),
            "sha256": sha256(data),
            "extension": path.suffix.lower(),
        }
        try:
            with Image.open(path) as image:
                info.update(
                    width=image.width,
                    height=image.height,
                    mode=image.mode,
                    has_alpha="A" in image.getbands(),
                )
        except Exception as error:
            info["image_error"] = type(error).__name__
        media[info["sha256"]] = info

    presentation = Presentation(args.pptx)
    assert len(presentation.slides) == 54, "unexpected slide count"
    slides = []
    used_by = defaultdict(list)
    for number, slide in enumerate(presentation.slides, start=1):
        stats = {
            "all_shapes": 0,
            "native_primitives": 0,
            "groups": 0,
            "pictures": 0,
            "connectors": 0,
            "tables": 0,
            "charts": 0,
            "text_shapes": 0,
            "text_chars": 0,
            "top_picture_area_ratio": 0.0,
            "max_depth": 0,
            "kinds": Counter(),
            "media_hashes": [],
        }
        records = analyze_shapes(
            slide.shapes,
            presentation.slide_width,
            presentation.slide_height,
            media,
            stats,
        )
        for digest in set(stats["media_hashes"]):
            used_by[digest].append(number)
        native_total = stats["native_primitives"] + stats["connectors"] + stats["tables"]
        leaf_total = native_total + stats["pictures"] + stats["charts"]
        slides.append(
            {
                "slide": number,
                "counts": {
                    **{key: value for key, value in stats.items() if key not in {"kinds", "media_hashes"}},
                    "kinds": dict(stats["kinds"]),
                    "native_leaf_ratio": round(native_total / leaf_total, 3) if leaf_total else 0,
                    "top_picture_area_ratio": round(stats["top_picture_area_ratio"], 3),
                    "unique_media": len(set(stats["media_hashes"])),
                },
                "shapes": records,
                "group_candidates": summarize_groups(records),
            }
        )

    for digest, item in media.items():
        item["used_by_slides"] = used_by.get(digest, [])

    output = {
        "source": str(args.pptx),
        "slide_size": {
            "width_emu": presentation.slide_width,
            "height_emu": presentation.slide_height,
        },
        "slides": slides,
        "media": sorted(media.values(), key=lambda item: item["name"]),
    }
    args.output.write_text(json.dumps(output, ensure_ascii=False, indent=2), encoding="utf-8")
    print(args.output)


if __name__ == "__main__":
    main()
